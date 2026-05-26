"""
Generate PowerPoint presentation for Enterprise HR RAG Platform
Run: python3 generate_presentation.py
Output: Enterprise_HR_RAG_Platform.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Brand Colours ──────────────────────────────────────────────────────────────
NAVY       = RGBColor(0x1B, 0x3A, 0x6B)   # primary dark
BLUE       = RGBColor(0x2E, 0x75, 0xB6)   # primary
LIGHT_BLUE = RGBColor(0xBD, 0xD7, 0xEE)   # accent light
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GREY  = RGBColor(0x26, 0x26, 0x26)
MID_GREY   = RGBColor(0x59, 0x59, 0x59)
LIGHT_GREY = RGBColor(0xF2, 0xF2, 0xF2)
GREEN      = RGBColor(0x1E, 0x6B, 0x3C)
AMBER      = RGBColor(0xC5, 0x5A, 0x11)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ── Helpers ────────────────────────────────────────────────────────────────────

def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    blank_layout = prs.slide_layouts[6]
    return prs.slides.add_slide(blank_layout)


def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, text, left, top, width, height,
                 font_size=18, bold=False, color=DARK_GREY,
                 align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_para(tf, text, font_size=14, bold=False, color=DARK_GREY,
             align=PP_ALIGN.LEFT, italic=False, space_before=0):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return p


def header_bar(slide, title, subtitle=None):
    """Dark navy header bar across the top."""
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.1), fill_color=NAVY)
    add_text_box(slide, title,
                 Inches(0.4), Inches(0.08), Inches(12), Inches(0.6),
                 font_size=28, bold=True, color=WHITE)
    if subtitle:
        add_text_box(slide, subtitle,
                     Inches(0.4), Inches(0.65), Inches(12), Inches(0.38),
                     font_size=14, color=LIGHT_BLUE)


def footer_bar(slide, text="ChandraAILabs · Enterprise HR RAG Platform · Confidential"):
    add_rect(slide, 0, SLIDE_H - Inches(0.32), SLIDE_W, Inches(0.32), fill_color=NAVY)
    add_text_box(slide, text,
                 Inches(0.3), SLIDE_H - Inches(0.3), Inches(12.5), Inches(0.28),
                 font_size=9, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)


def bullet_box(slide, items, left, top, width, height,
               font_size=14, color=DARK_GREY, bullet="•"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = f"{bullet}  {item}"
        run.font.size = Pt(font_size)
        run.font.color.rgb = color


def kv_table(slide, rows, left, top, width, col_widths=None,
             header_fill=NAVY, row_fill_alt=LIGHT_GREY, font_size=13):
    """Minimal table: first row = header (white text on navy)."""
    n_cols = len(rows[0])
    n_rows = len(rows)
    row_h  = Inches(0.38)
    table  = slide.shapes.add_table(n_rows, n_cols, left, top, width, row_h * n_rows).table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    for r_idx, row in enumerate(rows):
        for c_idx, cell_text in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = str(cell_text)
            p = cell.text_frame.paragraphs[0]
            run = p.runs[0] if p.runs else p.add_run()
            run.font.size = Pt(font_size)
            if r_idx == 0:
                run.font.bold = True
                run.font.color.rgb = WHITE
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_fill
            else:
                run.font.color.rgb = DARK_GREY
                if r_idx % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = row_fill_alt
                else:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = WHITE
    return table


# ══════════════════════════════════════════════════════════════════════════════
# SLIDES
# ══════════════════════════════════════════════════════════════════════════════

def slide_title(prs):
    slide = blank_slide(prs)
    # Full background gradient simulation via two rectangles
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_color=NAVY)
    add_rect(slide, 0, SLIDE_H * 0.55, SLIDE_W, SLIDE_H * 0.45, fill_color=BLUE)

    # Company tag
    add_text_box(slide, "CHANDRAAILABS",
                 Inches(0.5), Inches(0.35), Inches(12), Inches(0.5),
                 font_size=13, bold=True, color=LIGHT_BLUE,
                 align=PP_ALIGN.CENTER)

    # Main title
    add_text_box(slide, "Enterprise HR RAG Platform",
                 Inches(0.5), Inches(1.1), Inches(12.3), Inches(1.4),
                 font_size=44, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER)

    # Subtitle
    add_text_box(slide,
                 "Production-Grade Retrieval-Augmented Generation on Google Cloud Platform",
                 Inches(0.5), Inches(2.55), Inches(12.3), Inches(0.7),
                 font_size=18, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)

    # Divider line
    add_rect(slide, Inches(3.5), Inches(3.4), Inches(6.33), Inches(0.04), fill_color=LIGHT_BLUE)

    # Key tags
    tags = ["Hybrid RAG", "Google Cloud", "Gemini 2.5 Flash",
            "Vertex AI Vector Search", "Enterprise Security"]
    tag_x = Inches(1.0)
    for tag in tags:
        w = Inches(2.0)
        add_rect(slide, tag_x, Inches(3.65), w, Inches(0.38),
                 fill_color=BLUE, line_color=LIGHT_BLUE, line_width=Pt(1))
        add_text_box(slide, tag, tag_x + Inches(0.05), Inches(3.68),
                     w - Inches(0.1), Inches(0.34),
                     font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        tag_x += Inches(2.27)

    # Author block
    add_text_box(slide,
                 "Chandra Nakkalakunta  ·  Principal Cloud & AI/ML Architect  ·  Hyderabad, India",
                 Inches(0.5), Inches(4.5), Inches(12.3), Inches(0.4),
                 font_size=13, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide,
                 "linkedin.com/in/chandra-nakkalakunta  ·  chandraailabs.com",
                 Inches(0.5), Inches(4.95), Inches(12.3), Inches(0.35),
                 font_size=11, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)

    # Live demo URL
    add_rect(slide, Inches(3.2), Inches(5.55), Inches(6.93), Inches(0.5),
             fill_color=GREEN)
    add_text_box(slide,
                 "🌐  Live Demo: hr-rag-engine-946703664996.asia-south1.run.app",
                 Inches(3.25), Inches(5.58), Inches(6.83), Inches(0.44),
                 font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


def slide_agenda(prs):
    slide = blank_slide(prs)
    header_bar(slide, "Agenda")
    footer_bar(slide)

    items = [
        ("1", "Problem Statement & Business Value"),
        ("2", "System Architecture"),
        ("3", "Hybrid RAG — Retrieval Engine"),
        ("4", "Personal RAG — Personalised Answers"),
        ("5", "Response Caching (L1 + L2)"),
        ("6", "Security Architecture"),
        ("7", "Privacy & Compliance"),
        ("8", "Analytics & User Feedback"),
        ("9", "Operational Monitoring"),
        ("10", "Performance & Evaluation"),
        ("11", "Scalability"),
        ("12", "Infrastructure & GCP Resources"),
        ("13", "CI/CD Pipeline"),
        ("14", "Roadmap"),
    ]
    mid = len(items) // 2
    for col, chunk in enumerate([items[:mid], items[mid:]]):
        x = Inches(0.4) + col * Inches(6.5)
        for i, (num, text) in enumerate(chunk):
            y = Inches(1.35) + i * Inches(0.4)
            add_rect(slide, x, y, Inches(0.38), Inches(0.32), fill_color=NAVY)
            add_text_box(slide, num, x + Inches(0.02), y + Inches(0.01),
                         Inches(0.34), Inches(0.3),
                         font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
            add_text_box(slide, text, x + Inches(0.48), y,
                         Inches(5.8), Inches(0.35),
                         font_size=13, color=DARK_GREY)


def slide_problem(prs):
    slide = blank_slide(prs)
    header_bar(slide, "Problem Statement & Business Value",
               "Why enterprises need smarter HR Q&A")
    footer_bar(slide)

    # Problems column
    add_rect(slide, Inches(0.3), Inches(1.25), Inches(5.9), Inches(5.85),
             fill_color=LIGHT_GREY)
    add_text_box(slide, "❌  Current Challenges",
                 Inches(0.45), Inches(1.35), Inches(5.7), Inches(0.4),
                 font_size=16, bold=True, color=NAVY)

    problems = [
        "HR team spends 30-40% of time on repetitive policy questions",
        "Employees wait hours for answers to leave / WFH / CTC queries",
        "Policy documents scattered across SharePoint, Drive, Confluence",
        "No unified personal data view — employee must switch 3-4 systems",
        "Compliance risk from informal/verbal HR guidance",
        "Zero analytics on what employees actually struggle with",
    ]
    bullet_box(slide, problems, Inches(0.45), Inches(1.85),
               Inches(5.7), Inches(4.9), font_size=13, color=MID_GREY)

    # Solutions column
    add_rect(slide, Inches(6.5), Inches(1.25), Inches(6.5), Inches(5.85),
             fill_color=RGBColor(0xE8, 0xF4, 0xF8))
    add_text_box(slide, "✅  How This Platform Solves It",
                 Inches(6.65), Inches(1.35), Inches(6.3), Inches(0.4),
                 font_size=16, bold=True, color=GREEN)

    solutions = [
        "Instant natural language answers from official HR policy docs",
        "Personal data (leaves, CTC, rating) served directly from HR DB",
        "Hybrid RAG retrieves from all document sources simultaneously",
        "Query router auto-detects intent — personal vs policy vs hybrid",
        "Every answer cites the authoritative source document",
        "BigQuery analytics reveal top pain points — PII-free",
    ]
    bullet_box(slide, solutions, Inches(6.65), Inches(1.85),
               Inches(6.3), Inches(4.9), font_size=13, color=MID_GREY)


def slide_architecture(prs):
    slide = blank_slide(prs)
    header_bar(slide, "System Architecture",
               "End-to-end flow from employee question to personalised answer")
    footer_bar(slide)

    # Flow boxes
    boxes = [
        (Inches(0.3),  Inches(1.3),  Inches(2.0), "Employee\n(Browser)", NAVY),
        (Inches(2.6),  Inches(1.3),  Inches(2.0), "Google\nOAuth 2.0", BLUE),
        (Inches(5.0),  Inches(1.3),  Inches(2.0), "Streamlit\nUI (Cloud Run)", BLUE),
        (Inches(7.4),  Inches(1.3),  Inches(2.0), "Query\nRouter", NAVY),
        (Inches(10.0), Inches(1.3),  Inches(2.0), "Gemini 2.5\nFlash", GREEN),
    ]
    for x, y, w, label, col in boxes:
        add_rect(slide, x, y, w, Inches(0.75), fill_color=col)
        add_text_box(slide, label, x + Inches(0.05), y + Inches(0.05),
                     w - Inches(0.1), Inches(0.65),
                     font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # Arrow labels between top row
    arrows = ["OAuth", "Query", "Intent", "Prompt+Chunks"]
    arrow_x = [Inches(2.35), Inches(4.75), Inches(7.2), Inches(9.67)]
    for ax, al in zip(arrow_x, arrows):
        add_text_box(slide, f"→ {al}", ax, Inches(1.52),
                     Inches(0.55), Inches(0.3),
                     font_size=9, color=MID_GREY, align=PP_ALIGN.CENTER)

    # Three retrieval paths
    path_data = [
        (Inches(0.5),  "Personal Intent", "Cloud SQL\nPostgreSQL\n(HR Database)", NAVY),
        (Inches(5.2),  "Policy Intent",   "BM25 + Vector\nSearch (RRF\nFusion)", BLUE),
        (Inches(9.8),  "Hybrid Intent",   "DB + Retrieval\nCombined", AMBER),
    ]
    for px, plabel, pbox, pcol in path_data:
        add_rect(slide, px, Inches(2.7), Inches(3.0), Inches(0.3), fill_color=pcol)
        add_text_box(slide, plabel, px + Inches(0.05), Inches(2.72),
                     Inches(2.9), Inches(0.26),
                     font_size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_rect(slide, px, Inches(3.1), Inches(3.0), Inches(0.8), fill_color=LIGHT_GREY)
        add_text_box(slide, pbox, px + Inches(0.05), Inches(3.12),
                     Inches(2.9), Inches(0.76),
                     font_size=11, color=NAVY, align=PP_ALIGN.CENTER)

    # Answer + storage row
    add_rect(slide, Inches(4.2), Inches(4.2), Inches(5.0), Inches(0.65), fill_color=GREEN)
    add_text_box(slide, "Personalised Answer + Source Citations",
                 Inches(4.3), Inches(4.28), Inches(4.8), Inches(0.5),
                 font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    stores = [
        (Inches(0.3),  "Firestore\nCache (L2)", BLUE),
        (Inches(2.7),  "BigQuery\nAnalytics", NAVY),
        (Inches(5.1),  "Cloud\nLogging", BLUE),
        (Inches(7.5),  "Cloud\nMonitoring", NAVY),
        (Inches(9.9),  "User\nFeedback BQ", GREEN),
    ]
    add_text_box(slide, "Observability & Storage Layer",
                 Inches(0.3), Inches(5.05), Inches(12.5), Inches(0.28),
                 font_size=10, bold=True, color=NAVY)
    for sx, slabel, scol in stores:
        add_rect(slide, sx, Inches(5.35), Inches(2.2), Inches(0.65), fill_color=scol)
        add_text_box(slide, slabel, sx + Inches(0.05), Inches(5.38),
                     Inches(2.1), Inches(0.59),
                     font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


def slide_hybrid_rag(prs):
    slide = blank_slide(prs)
    header_bar(slide, "Hybrid RAG — Retrieval Engine",
               "BM25 sparse + Vertex AI Vector Search dense, fused with Reciprocal Rank Fusion")
    footer_bar(slide)

    # Left panel: retrieval pipeline
    add_rect(slide, Inches(0.3), Inches(1.25), Inches(6.2), Inches(5.85),
             fill_color=LIGHT_GREY)
    add_text_box(slide, "Retrieval Pipeline",
                 Inches(0.45), Inches(1.35), Inches(6.0), Inches(0.35),
                 font_size=15, bold=True, color=NAVY)
    steps = [
        ("1. User Query", "Raw natural language question"),
        ("2. BM25 Sparse", "Keyword matching, top-k candidates from local index"),
        ("3. Gemini Embedding", "gemini-embedding-001 (3072 dims) → dense vector"),
        ("4. Vector Search", "Vertex AI ANN index (STREAM_UPDATE mode), asia-south1"),
        ("5. RRF Fusion", "Reciprocal Rank Fusion merges both result sets"),
        ("6. Top-3 Chunks", "Highest-scored passages sent to Gemini"),
        ("7. Generation", "Gemini 2.5 Flash generates cited answer"),
    ]
    for i, (step, desc) in enumerate(steps):
        y = Inches(1.82) + i * Inches(0.72)
        add_rect(slide, Inches(0.45), y, Inches(2.0), Inches(0.38), fill_color=NAVY)
        add_text_box(slide, step, Inches(0.48), y + Inches(0.04),
                     Inches(1.94), Inches(0.3),
                     font_size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text_box(slide, desc, Inches(2.6), y + Inches(0.04),
                     Inches(3.8), Inches(0.34),
                     font_size=11, color=MID_GREY)

    # Right panel: config and performance
    add_rect(slide, Inches(6.8), Inches(1.25), Inches(6.1), Inches(2.6),
             fill_color=RGBColor(0xE8, 0xF4, 0xF8))
    add_text_box(slide, "Configuration",
                 Inches(6.95), Inches(1.35), Inches(5.9), Inches(0.35),
                 font_size=14, bold=True, color=NAVY)
    config = [
        "chunk_size = 1024 tokens",
        "chunk_overlap = 200 tokens",
        "alpha = 0.2 (BM25 weight)",
        "top_k_final = 3 chunks to LLM",
        "index dims = 3072",
        "distance_measure = DOT_PRODUCT",
    ]
    bullet_box(slide, config, Inches(6.95), Inches(1.78),
               Inches(5.9), Inches(1.9), font_size=12)

    add_rect(slide, Inches(6.8), Inches(4.0), Inches(6.1), Inches(3.1),
             fill_color=LIGHT_GREY)
    add_text_box(slide, "Evaluation Results (60 questions · RAGAS)",
                 Inches(6.95), Inches(4.1), Inches(5.9), Inches(0.35),
                 font_size=14, bold=True, color=NAVY)
    rows = [
        ["Metric", "BM25 Only", "Hybrid RAG", "Latest"],
        ["Avg Relevancy", "0.635", "0.674", "0.863"],
        ["Source Accuracy", "1.000", "1.000", "0.967"],
        ["Easy Questions", "0.706", "0.740", "0.850"],
        ["Medium Questions", "0.611", "0.630", "0.910"],
        ["Hard Questions", "0.377", "0.522", "0.830"],
    ]
    kv_table(slide, rows, Inches(6.95), Inches(4.5), Inches(5.8),
             col_widths=[Inches(2.0), Inches(1.2), Inches(1.3), Inches(1.1)],
             font_size=11)


def slide_personal_rag(prs):
    slide = blank_slide(prs)
    header_bar(slide, "Personal RAG — Personalised Answers",
               "Employee-specific data from Cloud SQL merged with policy context")
    footer_bar(slide)

    rows = [
        ["Question", "Intent", "Data Source", "Example Answer"],
        ["How many leaves do I have?", "personal", "Cloud SQL HR DB", "You have 12 casual + 0 sick leaves remaining"],
        ["What is my CTC?", "personal", "Cloud SQL HR DB", "Your CTC is ₹18.0L as of Jan 2024"],
        ["What is my performance rating?", "personal", "Cloud SQL HR DB", "Your latest rating: 4.5 / 5.0 (Exceeds)"],
        ["Who is my manager?", "personal", "Cloud SQL HR DB", "Your reporting manager is Rajesh Kumar"],
        ["Am I eligible for WFH?", "hybrid", "DB + Policy Docs", "Yes — Senior Engineers can WFH 3 days/week"],
        ["What is the WFH policy?", "policy", "Policy Docs (RAG)", "According to Remote Work Policy v2.1..."],
    ]
    kv_table(slide, rows, Inches(0.3), Inches(1.3), Inches(12.7),
             col_widths=[Inches(3.2), Inches(1.3), Inches(2.4), Inches(5.6)],
             font_size=11)

    add_text_box(slide, "Query Router — Intent Detection",
                 Inches(0.3), Inches(4.55), Inches(12.5), Inches(0.35),
                 font_size=14, bold=True, color=NAVY)

    intents = [
        ("Personal", "Regex detects \\b(my|i)\\b + personal data keywords\n→ Direct Cloud SQL query, no vector search", NAVY),
        ("Hybrid", "Personal + policy context needed\n→ DB lookup + RAG retrieval + combined prompt", BLUE),
        ("Policy", "Policy-only keywords (what is, how to, WFH policy)\n→ BM25 + Vector Search + Gemini generation", GREEN),
    ]
    for i, (title, desc, col) in enumerate(intents):
        x = Inches(0.3) + i * Inches(4.35)
        add_rect(slide, x, Inches(5.0), Inches(4.1), Inches(0.35), fill_color=col)
        add_text_box(slide, title, x + Inches(0.05), Inches(5.03),
                     Inches(4.0), Inches(0.29),
                     font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_rect(slide, x, Inches(5.38), Inches(4.1), Inches(1.5), fill_color=LIGHT_GREY)
        add_text_box(slide, desc, x + Inches(0.1), Inches(5.42),
                     Inches(3.9), Inches(1.4),
                     font_size=11, color=MID_GREY)


def slide_caching(prs):
    slide = blank_slide(prs)
    header_bar(slide, "Response Caching — L1 + L2",
               "Two-level cache for sub-second policy responses · Personal queries never cached")
    footer_bar(slide)

    levels = [
        (Inches(0.3), "L1 — In-Memory Cache",
         ["Type: Python dict per Cloud Run instance",
          "TTL: 30 minutes",
          "Size: 100 entries (LRU eviction)",
          "Scope: Single container instance",
          "BQ tag: model_used = cache_l1",
          "Latency on hit: ~0 ms"],
         NAVY, "Fastest — no I/O"),
        (Inches(4.55), "L2 — Firestore Cache",
         ["Type: Firestore 'rag_cache' collection",
          "TTL: 24 hours (expires_at field)",
          "Size: Unlimited (GCP managed)",
          "Scope: Shared across all instances",
          "BQ tag: model_used = cache_l2",
          "Latency on hit: ~30-50 ms"],
         BLUE, "Survives redeployments"),
        (Inches(8.8), "Cache Miss — Full Pipeline",
         ["BM25 keyword search",
          "Vertex AI Vector Search",
          "RRF fusion",
          "Gemini 2.5 Flash generation",
          "Write to L1 + L2 (non-blocking)",
          "BQ tag: model_used = gemini-2.5-flash"],
         AMBER, "~3-5 seconds"),
    ]
    for x, title, items, col, note in levels:
        add_rect(slide, x, Inches(1.25), Inches(4.0), Inches(0.38), fill_color=col)
        add_text_box(slide, title, x + Inches(0.08), Inches(1.28),
                     Inches(3.84), Inches(0.32),
                     font_size=13, bold=True, color=WHITE)
        add_rect(slide, x, Inches(1.65), Inches(4.0), Inches(4.0), fill_color=LIGHT_GREY)
        bullet_box(slide, items, x + Inches(0.1), Inches(1.72),
                   Inches(3.8), Inches(3.8), font_size=12)
        add_rect(slide, x, Inches(5.75), Inches(4.0), Inches(0.32), fill_color=col)
        add_text_box(slide, note, x + Inches(0.08), Inches(5.77),
                     Inches(3.84), Inches(0.28),
                     font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_text_box(slide,
                 "⚠️  Personal & hybrid queries are NEVER cached — employee data must not be shared across users",
                 Inches(0.3), Inches(6.2), Inches(12.7), Inches(0.36),
                 font_size=12, bold=True, color=AMBER)


def slide_security(prs):
    slide = blank_slide(prs)
    header_bar(slide, "Security Architecture",
               "Defence-in-depth across identity, container, network, data, and audit layers")
    footer_bar(slide)

    rows = [
        ["Layer", "Control", "Implementation"],
        ["Identity & Auth", "Google OAuth 2.0", "Company email domain enforcement; session tokens in-memory only"],
        ["Container Security", "Binary Authorization", "Only KMS RSA-4096 signed images permitted on Cloud Run"],
        ["Secret Management", "Secret Manager", "GEMINI_API_KEY, DB password — never in env vars or code"],
        ["Service Identity", "4 Service Accounts", "rag-engine-sa, ingestion-sa, evaluation-sa, cicd-sa (least privilege)"],
        ["Data Encryption", "CMEK + Cloud KMS", "hr-rag-keyring encrypts GCS buckets, Cloud SQL, Firestore"],
        ["Network", "VPC + Private Subnets", "Cloud SQL on private IP; Cloud NAT for outbound; no public DB"],
        ["Audit", "Cloud Logging", "All API calls logged; structured JSON with severity, user_id hash"],
        ["Container Image", "Artifact Registry", "Signed images only; vulnerability scanning on push"],
    ]
    kv_table(slide, rows, Inches(0.3), Inches(1.3), Inches(12.7),
             col_widths=[Inches(2.3), Inches(2.3), Inches(8.0)],
             font_size=11)

    add_text_box(slide, "Zero-Trust Principle: every component authenticates with a dedicated service account; no shared credentials",
                 Inches(0.3), Inches(6.3), Inches(12.7), Inches(0.3),
                 font_size=11, italic=True, color=MID_GREY)


def slide_privacy(prs):
    slide = blank_slide(prs)
    header_bar(slide, "Privacy & Compliance",
               "GDPR Article 5 · India DPDP Act 2023 · Zero PII in analytics")
    footer_bar(slide)

    # Left: what IS stored
    add_rect(slide, Inches(0.3), Inches(1.25), Inches(5.9), Inches(5.35),
             fill_color=RGBColor(0xE8, 0xF8, 0xE8))
    add_text_box(slide, "✅  BigQuery STORES (anonymised)",
                 Inches(0.45), Inches(1.35), Inches(5.7), Inches(0.35),
                 font_size=14, bold=True, color=GREEN)
    stored = [
        "hashed_user_id — SHA-256, irreversible, 16-char prefix",
        "question_category — leave / wfh / compensation / ...",
        "intent — personal / policy / hybrid",
        "department — Engineering / HR / Finance / ...",
        "latency_ms — response time",
        "model_used — cache_l1 / cache_l2 / gemini-2.5-flash",
        "success — true / false",
        "session_id — UUID per browser session",
        "timestamp — UTC (BQ partition key)",
        "feedback — positive / negative (feedback_logs table)",
    ]
    bullet_box(slide, stored, Inches(0.45), Inches(1.78), Inches(5.7), Inches(4.5),
               font_size=12, color=MID_GREY)

    # Right: what is NEVER stored
    add_rect(slide, Inches(6.5), Inches(1.25), Inches(6.5), Inches(2.8),
             fill_color=RGBColor(0xFF, 0xEB, 0xEB))
    add_text_box(slide, "❌  BigQuery NEVER STORES",
                 Inches(6.65), Inches(1.35), Inches(6.3), Inches(0.35),
                 font_size=14, bold=True, color=AMBER)
    never = [
        "Employee name or email address",
        "Actual question text",
        "Actual answer text",
        "CTC or salary figures",
        "Leave balance",
        "Performance ratings",
        "Manager names",
    ]
    bullet_box(slide, never, Inches(6.65), Inches(1.78), Inches(6.3), Inches(2.2),
               font_size=12, color=MID_GREY)

    # Compliance
    add_rect(slide, Inches(6.5), Inches(4.2), Inches(6.5), Inches(2.4),
             fill_color=LIGHT_GREY)
    add_text_box(slide, "Compliance Mapping",
                 Inches(6.65), Inches(4.3), Inches(6.3), Inches(0.35),
                 font_size=14, bold=True, color=NAVY)
    rows = [
        ["Regulation", "How Met"],
        ["GDPR Art 5", "Data minimisation — no PII"],
        ["India DPDP 2023", "Anonymised analytics only"],
        ["Right to Erasure", "Firestore TTL + partitioned BQ"],
        ["Data Retention", "BQ partitioned → auto-expire"],
    ]
    kv_table(slide, rows, Inches(6.65), Inches(4.7), Inches(6.1),
             col_widths=[Inches(2.2), Inches(3.8)], font_size=11)


def slide_analytics(prs):
    slide = blank_slide(prs)
    header_bar(slide, "Analytics & User Feedback",
               "Live BigQuery data · PII-free · Actionable insights for HR and Engineering")
    footer_bar(slide)

    # Usage data table
    add_text_box(slide, "Query Volume by Category (live data — 2026-05-26)",
                 Inches(0.3), Inches(1.28), Inches(12.5), Inches(0.33),
                 font_size=13, bold=True, color=NAVY)
    rows = [
        ["Category", "Intent", "Queries", "Avg Latency", "Cache Hits"],
        ["leave",       "policy",   "209", "4,351 ms", "9 (4%)"],
        ["wfh",         "policy",   "101", "4,536 ms", "3 (3%)"],
        ["performance", "policy",    "51", "4,793 ms", "1 (2%)"],
        ["travel",      "policy",    "41", "3,535 ms", "0"],
        ["performance", "hybrid",    "30", "5,400 ms", "0"],
        ["general",     "personal",  "23", "2,195 ms", "0"],
        ["compensation","personal",  "15", "2,223 ms", "0"],
    ]
    kv_table(slide, rows, Inches(0.3), Inches(1.65), Inches(7.0),
             col_widths=[Inches(1.6), Inches(1.2), Inches(1.1), Inches(1.5), Inches(1.4)],
             font_size=11)

    # Key insights
    add_rect(slide, Inches(7.4), Inches(1.25), Inches(5.6), Inches(3.5),
             fill_color=LIGHT_GREY)
    add_text_box(slide, "Key Observations",
                 Inches(7.55), Inches(1.35), Inches(5.4), Inches(0.33),
                 font_size=13, bold=True, color=NAVY)
    insights = [
        "Personal queries (2-2.2s) are ~2× faster than policy — no vector search",
        "Hybrid queries (5.4s) slowest — DB + retrieval + generation",
        "Leave and WFH are most queried — useful for HR planning",
        "Cache hit rate grows with usage (survives via Firestore L2)",
    ]
    bullet_box(slide, insights, Inches(7.55), Inches(1.75), Inches(5.4), Inches(2.7),
               font_size=12, color=MID_GREY)

    # Feedback section
    add_rect(slide, Inches(7.4), Inches(4.9), Inches(5.6), Inches(2.2),
             fill_color=RGBColor(0xE8, 0xF4, 0xF8))
    add_text_box(slide, "User Feedback (👍 / 👎)",
                 Inches(7.55), Inches(5.0), Inches(5.4), Inches(0.33),
                 font_size=13, bold=True, color=NAVY)
    fb_items = [
        "Per-response thumbs up / thumbs down",
        "Logged to BigQuery feedback_logs table",
        "Anonymised — SHA-256 hashed user ID",
        "Enables satisfaction tracking by category",
        "Powers continuous improvement roadmap",
    ]
    bullet_box(slide, fb_items, Inches(7.55), Inches(5.38), Inches(5.4), Inches(1.6),
               font_size=12, color=MID_GREY)

    add_text_box(slide, "Sample BQ Query →",
                 Inches(0.3), Inches(5.05), Inches(7.0), Inches(0.3),
                 font_size=12, bold=True, color=NAVY)
    code = ("SELECT question_category, intent,\n"
            "  COUNTIF(feedback='positive') AS up,\n"
            "  ROUND(COUNTIF(feedback='positive')*100.0/COUNT(*),1) AS pct\n"
            "FROM `hr-rag-dev.hr_rag_metrics.feedback_logs`\n"
            "GROUP BY 1,2 ORDER BY up DESC")
    add_rect(slide, Inches(0.3), Inches(5.38), Inches(7.0), Inches(1.65),
             fill_color=RGBColor(0x1E, 0x1E, 0x1E))
    add_text_box(slide, code, Inches(0.4), Inches(5.42), Inches(6.8), Inches(1.55),
                 font_size=9.5, color=RGBColor(0x9C, 0xDC, 0xFE))


def slide_monitoring(prs):
    slide = blank_slide(prs)
    header_bar(slide, "Operational Monitoring",
               "Structured JSON logging + Cloud Monitoring dashboard + 4 alert policies")
    footer_bar(slide)

    # App layer
    add_rect(slide, Inches(0.3), Inches(1.25), Inches(5.9), Inches(3.1),
             fill_color=LIGHT_GREY)
    add_text_box(slide, "App Layer — src/monitoring/",
                 Inches(0.45), Inches(1.35), Inches(5.7), Inches(0.33),
                 font_size=13, bold=True, color=NAVY)
    app_items = [
        "Structured JSON logs — severity, message, timestamp, extra fields",
        "Every query logs: intent, latency_ms, chunks_retrieved, success",
        "Cloud Run stdout → Cloud Logging parses JSON automatically",
        "Per-user analytics to BigQuery (anonymised, non-blocking thread)",
    ]
    bullet_box(slide, app_items, Inches(0.45), Inches(1.75), Inches(5.7), Inches(2.5),
               font_size=12, color=MID_GREY)

    # GCP layer
    add_rect(slide, Inches(0.3), Inches(4.5), Inches(5.9), Inches(2.6),
             fill_color=RGBColor(0xE8, 0xF4, 0xF8))
    add_text_box(slide, "GCP Layer — Cloud Monitoring",
                 Inches(0.45), Inches(4.6), Inches(5.7), Inches(0.33),
                 font_size=13, bold=True, color=NAVY)
    gcp_items = [
        "8-widget dashboard (request rate, p50/p95/p99 latency, error rate, CPU, memory, instances)",
        "Uptime check — 5-min interval, asia-pacific probe",
    ]
    bullet_box(slide, gcp_items, Inches(0.45), Inches(5.0), Inches(5.7), Inches(1.9),
               font_size=12, color=MID_GREY)

    # Alert policies
    add_rect(slide, Inches(6.5), Inches(1.25), Inches(6.5), Inches(5.85),
             fill_color=LIGHT_GREY)
    add_text_box(slide, "Alert Policies (all → email notification)",
                 Inches(6.65), Inches(1.35), Inches(6.3), Inches(0.33),
                 font_size=13, bold=True, color=NAVY)
    alerts = [
        ("5xx Error Rate > 5%", "5 min sustained · triggers incident response"),
        ("p99 Latency > 5s", "5 min sustained · investigate retrieval perf"),
        ("Service Down", "Zero active instances for 10 min · critical"),
        ("Memory > 85%", "5 min sustained · scale up or OOM risk"),
    ]
    for i, (title, desc) in enumerate(alerts):
        y = Inches(1.82) + i * Inches(1.25)
        add_rect(slide, Inches(6.65), y, Inches(6.1), Inches(0.38), fill_color=AMBER)
        add_text_box(slide, f"🔔  {title}",
                     Inches(6.75), y + Inches(0.04), Inches(5.9), Inches(0.3),
                     font_size=12, bold=True, color=WHITE)
        add_rect(slide, Inches(6.65), y + Inches(0.38), Inches(6.1), Inches(0.72),
                 fill_color=WHITE)
        add_text_box(slide, desc,
                     Inches(6.75), y + Inches(0.42), Inches(5.9), Inches(0.64),
                     font_size=11, color=MID_GREY)


def slide_performance(prs):
    slide = blank_slide(prs)
    header_bar(slide, "Performance & Scalability",
               "Latency profile · cache impact · Cloud Run auto-scaling")
    footer_bar(slide)

    # Latency breakdown
    add_rect(slide, Inches(0.3), Inches(1.25), Inches(6.2), Inches(5.85),
             fill_color=LIGHT_GREY)
    add_text_box(slide, "Latency Profile",
                 Inches(0.45), Inches(1.35), Inches(6.0), Inches(0.33),
                 font_size=14, bold=True, color=NAVY)
    rows = [
        ["Path", "Typical Latency", "Notes"],
        ["L1 Cache hit", "~0 ms", "In-memory, no I/O"],
        ["L2 Cache hit", "30-50 ms", "Firestore read"],
        ["Personal query", "1.5-2.5 s", "SQL only, no retrieval"],
        ["Policy query (miss)", "3-5 s", "BM25 + VectorSearch + LLM"],
        ["Hybrid query (miss)", "4-6 s", "DB + retrieval + LLM"],
        ["Cold start", "8-12 s", "First request on new instance"],
    ]
    kv_table(slide, rows, Inches(0.45), Inches(1.75), Inches(5.9),
             col_widths=[Inches(2.2), Inches(1.6), Inches(2.0)], font_size=12)

    # Scalability
    add_rect(slide, Inches(6.8), Inches(1.25), Inches(6.1), Inches(2.8),
             fill_color=RGBColor(0xE8, 0xF4, 0xF8))
    add_text_box(slide, "Cloud Run Auto-Scaling",
                 Inches(6.95), Inches(1.35), Inches(5.9), Inches(0.33),
                 font_size=14, bold=True, color=NAVY)
    scale_items = [
        "Scale to zero when idle (cost-efficient)",
        "Scales out on concurrent requests",
        "Max instances configurable (default: 10)",
        "Concurrency = 80 requests per instance",
        "Min instances = 1 (avoids cold starts in prod)",
    ]
    bullet_box(slide, scale_items, Inches(6.95), Inches(1.75), Inches(5.9), Inches(2.2),
               font_size=12, color=MID_GREY)

    add_rect(slide, Inches(6.8), Inches(4.2), Inches(6.1), Inches(2.9),
             fill_color=LIGHT_GREY)
    add_text_box(slide, "Bottlenecks & Mitigations",
                 Inches(6.95), Inches(4.3), Inches(5.9), Inches(0.33),
                 font_size=14, bold=True, color=NAVY)
    bottlenecks = [
        "Vector Search cold index → mitigated by L1/L2 cache",
        "Gemini API latency → mitigated by streaming + cache",
        "Cloud SQL connection pool → pgbouncer in prod",
        "Streamlit single-thread → planned Chainlit migration",
        "BQ analytics write → daemon thread, never blocks",
    ]
    bullet_box(slide, bottlenecks, Inches(6.95), Inches(4.7), Inches(5.9), Inches(2.3),
               font_size=12, color=MID_GREY)


def slide_infrastructure(prs):
    slide = blank_slide(prs)
    header_bar(slide, "Infrastructure & GCP Resources",
               "Fully automated — 14 idempotent shell scripts + Makefile")
    footer_bar(slide)

    rows = [
        ["Resource", "Details"],
        ["Project ID", "hr-rag-dev | Region: asia-south1 (Mumbai)"],
        ["VPC Network", "hr-rag-vpc | 10.1.0.0/24 (app) + 10.2.0.0/24 (data)"],
        ["Cloud Run", "hr-rag-engine | auto-scaling 0-10 | concurrency 80"],
        ["Vertex AI Vector Search", "3072 dims · STREAM_UPDATE · DOT_PRODUCT distance"],
        ["Cloud SQL", "hr-rag-db · PostgreSQL · private IP · 12 demo employees"],
        ["Firestore", "rag_cache collection (L2) + hr_data (10 docs, 44 chunks)"],
        ["GCS Buckets", "documents / processed / artifacts / audit"],
        ["BigQuery", "hr_rag_metrics.query_logs + feedback_logs"],
        ["Secret Manager", "gemini-api-key + hr-db-password"],
        ["Cloud KMS", "hr-rag-keyring · RSA-4096 signing key"],
        ["Service Accounts", "rag-engine-sa / ingestion-sa / evaluation-sa / cicd-sa"],
        ["Cloud Monitoring", "Dashboard (8 widgets) + 4 alert policies + uptime check"],
        ["Cloud Build", "CI/CD pipeline · asia-south1 · Binary Auth signing"],
        ["APIs Enabled", "19 GCP APIs (Vertex AI, Cloud Run, KMS, IAM, BQ, Firestore…)"],
    ]
    kv_table(slide, rows, Inches(0.3), Inches(1.3), Inches(12.7),
             col_widths=[Inches(2.8), Inches(9.8)], font_size=11)


def slide_cicd(prs):
    slide = blank_slide(prs)
    header_bar(slide, "CI/CD Pipeline",
               "Cloud Build · Binary Authorization · KMS signing · Automated deployment")
    footer_bar(slide)

    steps = [
        ("1. git push", "Developer pushes to main branch"),
        ("2. Cloud Build Trigger", "Auto-triggered on push to main"),
        ("3. docker build", "Container built in asia-south1 region"),
        ("4. Push to Artifact Registry", "Image tagged with commit SHA"),
        ("5. KMS Sign Image", "RSA-4096 signature via Cloud KMS"),
        ("6. Binary Auth Check", "Cloud Run verifies signature before deploy"),
        ("7. Cloud Run Deploy", "Zero-downtime rolling deployment"),
        ("8. Smoke Test", "Health check + BQ query verify"),
    ]
    for i, (step, desc) in enumerate(steps):
        col = i % 2
        row = i // 2
        x = Inches(0.3) + col * Inches(6.5)
        y = Inches(1.3) + row * Inches(1.35)
        add_rect(slide, x, y, Inches(6.0), Inches(0.38), fill_color=NAVY)
        add_text_box(slide, step, x + Inches(0.1), y + Inches(0.04),
                     Inches(5.8), Inches(0.3),
                     font_size=12, bold=True, color=WHITE)
        add_rect(slide, x, y + Inches(0.38), Inches(6.0), Inches(0.82),
                 fill_color=LIGHT_GREY)
        add_text_box(slide, desc, x + Inches(0.1), y + Inches(0.44),
                     Inches(5.8), Inches(0.7),
                     font_size=12, color=MID_GREY)

    add_text_box(slide,
                 "Security gate: Binary Authorization BLOCKS deployment if image is unsigned — no exceptions",
                 Inches(0.3), Inches(6.82), Inches(12.7), Inches(0.3),
                 font_size=11, bold=True, color=AMBER)


def slide_roadmap(prs):
    slide = blank_slide(prs)
    header_bar(slide, "Roadmap",
               "What's done · What's next")
    footer_bar(slide)

    phases = [
        ("Phase 1 — Infrastructure", COMPLETE, [
            "VPC + subnets + Cloud NAT",
            "IAM + KMS + Binary Auth",
            "GCS + BigQuery + Firestore",
            "Vertex AI Vector Search",
            "Cloud SQL PostgreSQL",
            "CI/CD with Cloud Build",
        ]),
        ("Phase 2 — Personal RAG + Auth", COMPLETE, [
            "Google OAuth (company email)",
            "HR database with employee data",
            "Query router (personal/hybrid/policy)",
            "Personal RAG (leave, CTC, rating)",
            "Streamlit chat UI",
            "10 HR policy documents ingested",
        ]),
        ("Phase 3 — MLOps", PARTIAL, [
            "✅ Structured logging + Cloud Monitoring",
            "✅ L1 + L2 response caching",
            "✅ User feedback (👍/👎) to BigQuery",
            "✅ Query router fix (word boundary)",
            "🔲 Vertex AI Experiments tracking",
            "🔲 A/B testing framework",
        ]),
        ("Phase 4 — Advanced", PLANNED, [
            "Multi-tenancy (multiple orgs)",
            "Slack / Teams bot integration",
            "Chainlit / FastAPI+React UI",
            "Mobile app",
            "Drift monitoring",
            "Automated retraining",
        ]),
    ]
    return phases  # used below


COMPLETE = GREEN
PARTIAL  = BLUE
PLANNED  = MID_GREY


def slide_roadmap_real(prs):
    slide = blank_slide(prs)
    header_bar(slide, "Roadmap", "What's done · What's next")
    footer_bar(slide)

    phases = [
        ("Phase 1 — Infrastructure", "COMPLETE", GREEN, [
            "VPC + subnets + Cloud NAT",
            "IAM + KMS + Binary Auth",
            "GCS + BigQuery + Firestore",
            "Vertex AI Vector Search + Cloud SQL",
            "CI/CD with Cloud Build",
        ]),
        ("Phase 2 — Personal RAG + Auth", "COMPLETE", GREEN, [
            "Google OAuth (company email)",
            "HR database with employee data",
            "Query router (personal/hybrid/policy)",
            "Personal RAG + Streamlit chat UI",
            "10 HR policy documents ingested",
        ]),
        ("Phase 3 — MLOps", "IN PROGRESS", BLUE, [
            "✅ Structured JSON logging + Cloud Monitoring",
            "✅ L1 + L2 response caching (in-memory + Firestore)",
            "✅ User feedback (👍/👎) to BigQuery",
            "✅ Query router fix (word-boundary regex)",
            "🔲 Vertex AI Experiments tracking",
            "🔲 A/B testing framework",
        ]),
        ("Phase 4 — Advanced", "PLANNED", MID_GREY, [
            "Multi-tenancy support (multiple orgs)",
            "Slack / Teams bot integration",
            "Chainlit / FastAPI+React UI migration",
            "Mobile app",
            "Drift monitoring + automated retraining",
        ]),
    ]
    for i, (title, status, col, items) in enumerate(phases):
        x = Inches(0.3) + i * Inches(3.25)
        add_rect(slide, x, Inches(1.25), Inches(3.0), Inches(0.38), fill_color=col)
        add_text_box(slide, title, x + Inches(0.06), Inches(1.28),
                     Inches(2.88), Inches(0.3),
                     font_size=10, bold=True, color=WHITE)
        add_rect(slide, x, Inches(1.63), Inches(3.0), Inches(0.28), fill_color=LIGHT_GREY)
        add_text_box(slide, status, x + Inches(0.06), Inches(1.65),
                     Inches(2.88), Inches(0.24),
                     font_size=9, bold=True, color=col, align=PP_ALIGN.CENTER)
        add_rect(slide, x, Inches(1.93), Inches(3.0), Inches(4.8), fill_color=LIGHT_GREY)
        bullet_box(slide, items, x + Inches(0.1), Inches(2.0),
                   Inches(2.8), Inches(4.6), font_size=10, color=MID_GREY)


def slide_closing(prs):
    slide = blank_slide(prs)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill_color=NAVY)
    add_rect(slide, 0, SLIDE_H * 0.6, SLIDE_W, SLIDE_H * 0.4, fill_color=BLUE)

    add_text_box(slide, "Thank You",
                 Inches(0.5), Inches(0.7), Inches(12.3), Inches(1.1),
                 font_size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_text_box(slide, "Enterprise HR RAG Platform",
                 Inches(0.5), Inches(1.85), Inches(12.3), Inches(0.55),
                 font_size=22, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)

    add_rect(slide, Inches(3.5), Inches(2.55), Inches(6.33), Inches(0.04),
             fill_color=LIGHT_BLUE)

    details = [
        ("Author",   "Chandra Nakkalakunta"),
        ("Role",     "Principal Cloud & AI/ML Architect"),
        ("Location", "Hyderabad, India"),
        ("LinkedIn", "linkedin.com/in/chandra-nakkalakunta"),
        ("GitHub",   "github.com/chandranakkalakunta"),
        ("Website",  "chandraailabs.com"),
    ]
    for i, (label, value) in enumerate(details):
        y = Inches(2.85) + i * Inches(0.42)
        add_text_box(slide, f"{label}:", Inches(3.2), y, Inches(1.4), Inches(0.38),
                     font_size=13, bold=True, color=LIGHT_BLUE, align=PP_ALIGN.RIGHT)
        add_text_box(slide, value, Inches(4.75), y, Inches(5.5), Inches(0.38),
                     font_size=13, color=WHITE)

    add_rect(slide, Inches(3.2), Inches(5.6), Inches(6.93), Inches(0.5),
             fill_color=GREEN)
    add_text_box(slide,
                 "🌐  Live Demo: hr-rag-engine-946703664996.asia-south1.run.app",
                 Inches(3.25), Inches(5.63), Inches(6.83), Inches(0.44),
                 font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_text_box(slide,
                 "demo@chandraailabs.com / demo123  ·  guest@chandraailabs.com / guest123",
                 Inches(0.5), Inches(6.25), Inches(12.3), Inches(0.35),
                 font_size=11, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# MAIN
# ══════════════════════════════════════════════════════════════════════════════

def main():
    prs = new_prs()

    slide_title(prs)
    slide_agenda(prs)
    slide_problem(prs)
    slide_architecture(prs)
    slide_hybrid_rag(prs)
    slide_personal_rag(prs)
    slide_caching(prs)
    slide_security(prs)
    slide_privacy(prs)
    slide_analytics(prs)
    slide_monitoring(prs)
    slide_performance(prs)
    slide_infrastructure(prs)
    slide_cicd(prs)
    slide_roadmap_real(prs)
    slide_closing(prs)

    out = "Enterprise_HR_RAG_Platform.pptx"
    prs.save(out)
    print(f"✅  Saved: {out}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
